from __future__ import annotations

import asyncio
import io
import mimetypes
import uuid
from pathlib import Path

from datetime import date

from fastapi import APIRouter, HTTPException, Query, Request, UploadFile
from fastapi.responses import Response
from PIL import Image
from pydantic import BaseModel
from sqlalchemy import func, select
from sqlalchemy import text as sql

from rag_app.api.audit import audit
from rag_app.api.auth import User, require_user
from rag_app.api.schemas import DocumentOut
from rag_app.config import settings
from rag_app.db.models import Document, DocumentStatus, DocumentTranslation, Segment
from rag_app.pipeline import ooxml
from rag_app.rag.memory.rls import apply_scope_guc

router = APIRouter(prefix="/api/documents", tags=["documents"], dependencies=[require_user])


def _owner_filter(stmt, user: User):
    """RBAC: admin видит всё; user — свои + документы dev-периода (owner NULL)."""
    if user.is_admin:
        return stmt
    return stmt.where((Document.owner_sub == user.sub) | (Document.owner_sub.is_(None)))


_PREVIEW_RENDER_WIDTH = 900


def _preview_source(doc: Document) -> tuple[str, str] | None:
    """PDF-источник для точного превью первой страницы."""
    if doc.s3_key_view_orig:
        return settings.bucket_exports, doc.s3_key_view_orig
    if doc.content_type == "application/pdf" or doc.filename.lower().endswith(".pdf"):
        return settings.bucket_originals, doc.s3_key_original
    return None


def _render_pdf_preview_png(pdf_bytes: bytes) -> bytes:
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(pdf_bytes)
    try:
        if len(pdf) == 0:
            raise ValueError("empty PDF")
        page = pdf[0]
        try:
            width = max(float(page.get_width()), 1.0)
            scale = min(max(_PREVIEW_RENDER_WIDTH / width, 1.0), 2.5)
            image = page.render(scale=scale, draw_annots=True).to_pil().convert("RGB")
        finally:
            page.close()
    finally:
        pdf.close()
    out = io.BytesIO()
    image.save(out, format="PNG", optimize=True)
    return out.getvalue()

# ТЗ §4.2: PDF (текст/скан), OOXML, изображения документов (JPG/PNG → OCR-ветка),
# plain-текст (TXT). Изображения оборачиваются в 1-страничный PDF на загрузке.
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".jpg", ".jpeg", ".png", ".txt"}
_SIGNATURES = {
    ".pdf": b"%PDF",
    ".docx": b"PK",
    ".xlsx": b"PK",
    ".pptx": b"PK",
    ".jpg": b"\xff\xd8\xff",
    ".jpeg": b"\xff\xd8\xff",
    ".png": b"\x89PNG",
}
_IMAGE_EXTS = {".jpg", ".jpeg", ".png"}


@router.post("", response_model=DocumentOut, status_code=201)
async def upload_document(request: Request, file: UploadFile) -> DocumentOut:
    # Направление перевода НЕ выбирается вручную (ТЗ §4.3, домен ru/en/zh):
    # язык-источник определяется автоматически по тексту на этапе перевода,
    # цель — всегда русский (русский документ не переводится). "auto" —
    # маркер «ещё не определён», воркер заменит его на ru/en/zh.
    source_lang, target_lang = "auto", "ru"
    filename = file.filename or "document.pdf"
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        allowed = ", ".join(sorted(ALLOWED_EXTENSIONS))
        raise HTTPException(415, f"поддерживаются {allowed}, получено: {ext or 'без расширения'}")

    data = await file.read()
    if len(data) > settings.max_upload_mb * 1024 * 1024:
        raise HTTPException(413, f"файл больше {settings.max_upload_mb} МБ")
    sig = _SIGNATURES.get(ext)
    if sig and not data.startswith(sig):
        raise HTTPException(415, f"содержимое не похоже на {ext}")

    content_type = file.content_type or "application/octet-stream"
    # ТЗ §4.2: изображение документа/чертежа → 1-страничный PDF, дальше как скан
    if ext in _IMAGE_EXTS:
        try:
            img = Image.open(io.BytesIO(data)).convert("RGB")
            buf = io.BytesIO()
            img.save(buf, format="PDF", resolution=150.0)
        except Exception as exc:
            raise HTTPException(415, f"не удалось прочитать изображение: {exc}") from None
        data = buf.getvalue()
        filename = Path(filename).stem + ".pdf"
        content_type = "application/pdf"

    doc_id = uuid.uuid4()
    s3_key = f"{doc_id}/{filename}"
    await request.app.state.storage.put_bytes(
        settings.bucket_originals, s3_key, data, content_type=content_type
    )

    async with request.app.state.sessionmaker() as session:
        doc = Document(
            id=doc_id,
            owner_sub=request.state.user.sub if settings.auth_enabled else None,
            filename=filename,
            content_type=content_type,
            size_bytes=len(data),
            s3_key_original=s3_key,
            source_lang=source_lang,
            target_lang=target_lang,
        )
        session.add(doc)
        await session.commit()
        await session.refresh(doc)

    await request.app.state.arq.enqueue_job(
        "parse_document", str(doc_id), _job_id=f"parse:{doc_id}:{uuid.uuid4().hex[:8]}"
    )
    await audit(request, "upload", "document", str(doc_id), {"filename": filename, "bytes": len(data)})
    return DocumentOut.from_doc(doc)


@router.get("", response_model=list[DocumentOut])
async def list_documents(
    request: Request,
    kind: str | None = Query(None, description="тип файла: pdf_text|pdf_scan|docx|xlsx|pptx|text"),
    date_from: str | None = Query(None, description="дата загрузки от (YYYY-MM-DD)"),
    date_to: str | None = Query(None, description="дата загрузки до (YYYY-MM-DD)"),
) -> list[DocumentOut]:
    # Фильтры списка (ТЗ §4.7.3) — только тип и даты; поиск по имени/содержимому —
    # в едином /api/search (гибрид + reranker + имя файла).
    async with request.app.state.sessionmaker() as session:
        stmt = _owner_filter(select(Document), request.state.user)
        if kind and kind.strip():
            stmt = stmt.where(Document.kind == kind.strip())
        for raw, op in ((date_from, ">="), (date_to, "<=")):
            if raw:
                try:
                    d = date.fromisoformat(raw)
                except ValueError:
                    continue
                col = func.date(Document.created_at)
                stmt = stmt.where(col >= d if op == ">=" else col <= d)
        docs = (
            (await session.execute(stmt.order_by(Document.created_at.desc()).limit(200)))
            .scalars()
            .all()
        )
        reviews = dict(
            (
                await session.execute(
                    select(Segment.document_id, func.count())
                    .where(Segment.needs_review)
                    .group_by(Segment.document_id)
                )
            ).all()
        )
    return [DocumentOut.from_doc(d, reviews.get(d.id, 0)) for d in docs]


@router.get("/{doc_id}", response_model=DocumentOut)
async def get_document(request: Request, doc_id: uuid.UUID) -> DocumentOut:
    doc = await _get_or_404(request, doc_id)
    async with request.app.state.sessionmaker() as session:
        review_count = (
            await session.execute(
                select(func.count())
                .select_from(Segment)
                .where(Segment.document_id == doc_id, Segment.needs_review)
            )
        ).scalar_one()
    return DocumentOut.from_doc(doc, review_count)


@router.post("/{doc_id}/retry", response_model=DocumentOut)
async def retry_document(request: Request, doc_id: uuid.UUID) -> DocumentOut:
    """Перезапуск пайплайна с парсинга (после ошибки)."""
    doc = await _get_or_404(request, doc_id)
    if doc.status not in (DocumentStatus.error, DocumentStatus.done):
        raise HTTPException(409, f"документ в работе (статус {doc.status.value})")
    await request.app.state.arq.enqueue_job(
        "parse_document", str(doc_id), _job_id=f"parse:{doc_id}:{uuid.uuid4().hex[:8]}"
    )
    await audit(request, "retry", "document", str(doc_id))
    return DocumentOut.from_doc(doc)


@router.post("/{doc_id}/describe")
async def describe_document(request: Request, doc_id: uuid.UUID) -> dict:
    """Запуск VL-описания рисунков документа (скан/чертёж/P&ID) по требованию.
    Раскрывает смысл изображений текстом и переиндексирует (см. describe_images)."""
    doc = await _get_or_404(request, doc_id)
    if not settings.vl_enabled:
        raise HTTPException(409, "VL-описание выключено (vl_enabled=false)")
    await request.app.state.arq.enqueue_job(
        "describe_images", str(doc_id), _job_id=f"vl:{doc_id}:{uuid.uuid4().hex[:8]}"
    )
    await audit(request, "describe", "document", str(doc_id))
    return {"status": "queued", "kind": doc.kind}


class ReparseOcrIn(BaseModel):
    # en | east_slavic (рус/укр/бел) | cyrillic | latin | ch | … (см. mineru -l)
    lang: str = "east_slavic"


@router.post("/{doc_id}/reparse-ocr")
async def reparse_ocr(request: Request, doc_id: uuid.UUID, body: ReparseOcrIn) -> dict:
    """Переразбор через форс-OCR — восстановление PDF с битым ToUnicode-cmap
    текстового слоя (MinerU `-m ocr -l <lang>`); выбор сохраняется на документе."""
    async with request.app.state.sessionmaker() as session:
        doc = await session.get(Document, doc_id)
        if doc is None:
            raise HTTPException(404, "документ не найден")
        if doc.status not in (DocumentStatus.error, DocumentStatus.done):
            raise HTTPException(409, f"документ в работе (статус {doc.status.value})")
        doc.parse_force_ocr = True
        doc.ocr_lang = body.lang
        doc.status = DocumentStatus.uploaded
        await session.commit()
    await request.app.state.arq.enqueue_job(
        "parse_document", str(doc_id), _job_id=f"parse:{doc_id}:{uuid.uuid4().hex[:8]}"
    )
    await audit(request, "reparse_ocr", "document", str(doc_id), {"lang": body.lang})
    return {"status": "queued", "ocr_lang": body.lang}


_PARSER_BACKENDS = {"mineru", "dots_mocr", "paddle_vl"}


class ReparseIn(BaseModel):
    # mineru (MinerU2.5-Pro + добор) | dots_mocr | paddle_vl
    backend: str = "mineru"


@router.post("/{doc_id}/reparse")
async def reparse(request: Request, doc_id: uuid.UUID, body: ReparseIn) -> dict:
    """Переразбор выбранным движком парсинга pdf_text (mineru | dots_mocr |
    paddle_vl). Выбор сохраняется на документе и переживает retry/reexport."""
    if body.backend not in _PARSER_BACKENDS:
        raise HTTPException(422, f"неизвестный backend: {body.backend}")
    async with request.app.state.sessionmaker() as session:
        doc = await session.get(Document, doc_id)
        if doc is None:
            raise HTTPException(404, "документ не найден")
        if doc.status not in (DocumentStatus.error, DocumentStatus.done):
            raise HTTPException(409, f"документ в работе (статус {doc.status.value})")
        doc.parser_backend = body.backend
        doc.parse_force_ocr = False  # выбор движка и форс-OCR взаимоисключающи
        doc.status = DocumentStatus.uploaded
        await session.commit()
    await request.app.state.arq.enqueue_job(
        "parse_document", str(doc_id), _job_id=f"parse:{doc_id}:{uuid.uuid4().hex[:8]}"
    )
    await audit(request, "reparse", "document", str(doc_id), {"backend": body.backend})
    return {"status": "queued", "backend": body.backend}


_EXPORT_KINDS = {
    "docx": ("s3_key_export_docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"),
    "pdf": ("s3_key_export_pdf", "application/pdf"),
    "pdf_dual": ("s3_key_export_pdf_dual", "application/pdf"),
    "source": ("s3_key_export_source", "application/octet-stream"),
}


@router.get("/{doc_id}/download/{kind}")
async def download(request: Request, doc_id: uuid.UUID, kind: str) -> Response:
    doc = await _get_or_404(request, doc_id)
    if kind == "original":
        bucket, key = settings.bucket_originals, doc.s3_key_original
        media = doc.content_type or "application/octet-stream"
        out_name = doc.filename
    elif kind in ("view_orig", "view_ru"):
        # PDF-рендер OOXML для просмотра «как в Microsoft» (LibreOffice)
        key = doc.s3_key_view_orig if kind == "view_orig" else doc.s3_key_view_ru
        if not key:
            raise HTTPException(404, "PDF-просмотр не готов")
        bucket, media, out_name = settings.bucket_exports, "application/pdf", Path(key).name
    elif kind in _EXPORT_KINDS:
        attr, media = _EXPORT_KINDS[kind]
        key = getattr(doc, attr)
        if not key:
            raise HTTPException(404, f"экспорт «{kind}» ещё не готов")
        bucket = settings.bucket_exports
        out_name = Path(key).name
    else:
        raise HTTPException(404, f"kind должен быть один из: original, {', '.join(_EXPORT_KINDS)}")

    data = await request.app.state.storage.get_bytes(bucket, key)
    await audit(request, "download", "document", str(doc_id), {"kind": kind})
    quoted = out_name.encode("ascii", "ignore").decode() or "document"
    total = len(data)
    base = {
        "Content-Disposition": f'attachment; filename="{quoted}"',
        "Accept-Ranges": "bytes",  # pdf.js тянет страницы лениво порейндж-запросами
    }
    # HTTP Range (RFC 7233): большие PDF открываются сразу — pdf.js запрашивает
    # сначала структуру (хвост файла), затем страницы по мере листания, а не весь
    # файл целиком. Без Range отдаём всё, но с Accept-Ranges.
    rng = request.headers.get("range", "")
    if rng.startswith("bytes="):
        try:
            s, _, e = rng[6:].partition("-")
            start = int(s) if s else 0
            end = int(e) if e else total - 1
            end = min(end, total - 1)
            if start > end or start >= total:
                raise ValueError
        except ValueError:
            return Response(status_code=416, headers={"Content-Range": f"bytes */{total}", **base})
        chunk = data[start : end + 1]
        return Response(
            content=chunk,
            status_code=206,
            media_type=media,
            headers={**base, "Content-Range": f"bytes {start}-{end}/{total}", "Content-Length": str(len(chunk))},
        )
    return Response(content=data, media_type=media, headers={**base, "Content-Length": str(total)})


@router.get("/{doc_id}/preview.png")
async def document_preview(request: Request, doc_id: uuid.UUID) -> Response:
    doc = await _get_or_404(request, doc_id)
    source = _preview_source(doc)
    if source is None:
        raise HTTPException(404, "PDF-превью ещё не готово")

    cache_key = f"{doc_id}/preview/page-1.png"
    storage = request.app.state.storage
    try:
        data = await storage.get_bytes(settings.bucket_artifacts, cache_key)
        return Response(
            content=data,
            media_type="image/png",
            headers={"Cache-Control": "public, max-age=86400"},
        )
    except Exception:
        pass

    bucket, key = source
    try:
        pdf_bytes = await storage.get_bytes(bucket, key)
        data = await asyncio.to_thread(_render_pdf_preview_png, pdf_bytes)
    except Exception as exc:
        raise HTTPException(404, "не удалось подготовить превью первой страницы") from exc

    try:
        await storage.put_bytes(settings.bucket_artifacts, cache_key, data, "image/png")
    except Exception:
        pass
    return Response(
        content=data,
        media_type="image/png",
        headers={"Cache-Control": "public, max-age=86400"},
    )


# --- Доп. переводы документа (ТЗ §4.3): RU→EN / RU→ZH по запросу --------------

class TranslationIn(BaseModel):
    target_lang: str  # en | ru | zh


@router.post("/{doc_id}/translations", status_code=202)
async def create_translation(request: Request, doc_id: uuid.UUID, body: TranslationIn) -> dict:
    """Запустить перевод документа на дополнительный язык. Основной перевод (→ru)
    не трогается; результат — отдельный артефакт DocumentTranslation."""
    doc = await _get_or_404(request, doc_id)
    target = (body.target_lang or "").strip().lower()
    if target not in ("en", "ru", "zh"):
        raise HTTPException(422, "target_lang: en | ru | zh")
    src = (doc.source_lang or "").strip().lower() or "ru"
    if target == src:
        raise HTTPException(422, f"документ уже на языке «{target}» — переводить не нужно")
    if doc.status not in (DocumentStatus.done, DocumentStatus.translated, DocumentStatus.exporting):
        raise HTTPException(409, "документ ещё не обработан — дождитесь готовности")
    async with request.app.state.sessionmaker() as db:
        row = (
            await db.execute(
                select(DocumentTranslation).where(
                    DocumentTranslation.document_id == doc_id,
                    DocumentTranslation.target_lang == target,
                )
            )
        ).scalar_one_or_none()
        if row is None:
            row = DocumentTranslation(document_id=doc_id, target_lang=target, status="translating")
            db.add(row)
        else:  # перезапуск — сброс прошлого результата
            row.status = "translating"
            row.error = None
            row.data = {}
            row.translated_count = 0
            row.s3_key_docx = None
            row.s3_key_source = None
        await db.commit()
    await request.app.state.arq.enqueue_job(
        "translate_to_language", str(doc_id), target, _job_id=f"translate_lang:{doc_id}:{target}"
    )
    await audit(request, "translate_lang", "document", str(doc_id), {"target_lang": target})
    return {"target_lang": target, "status": "translating"}


@router.get("/{doc_id}/translations")
async def list_translations(request: Request, doc_id: uuid.UUID) -> list[dict]:
    await _get_or_404(request, doc_id)
    async with request.app.state.sessionmaker() as db:
        rows = (
            (
                await db.execute(
                    select(DocumentTranslation).where(DocumentTranslation.document_id == doc_id)
                )
            )
            .scalars()
            .all()
        )
    return [
        {
            "target_lang": r.target_lang,
            "status": r.status,
            "translated_count": r.translated_count,
            "segment_count": r.segment_count,
            "needs_review_count": r.needs_review_count,
            "has_export": bool(r.s3_key_docx or r.s3_key_source),
            "error": r.error,
        }
        for r in rows
    ]


@router.get("/{doc_id}/translations/{lang}/download")
async def download_translation(request: Request, doc_id: uuid.UUID, lang: str) -> Response:
    await _get_or_404(request, doc_id)
    async with request.app.state.sessionmaker() as db:
        row = (
            await db.execute(
                select(DocumentTranslation).where(
                    DocumentTranslation.document_id == doc_id,
                    DocumentTranslation.target_lang == lang,
                )
            )
        ).scalar_one_or_none()
    key = (row.s3_key_docx or row.s3_key_source) if row else None
    if not key:
        raise HTTPException(404, "перевод ещё не готов")
    media = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        if row.s3_key_docx
        else "application/octet-stream"
    )
    data = await request.app.state.storage.get_bytes(settings.bucket_exports, key)
    await audit(request, "download_translation", "document", str(doc_id), {"lang": lang})
    name = Path(key).name.encode("ascii", "ignore").decode() or f"translation.{lang}"
    return Response(
        content=data, media_type=media,
        headers={"Content-Disposition": f'attachment; filename="{name}"'},
    )


@router.get("/{doc_id}/image/{name}")
async def document_image(request: Request, doc_id: uuid.UUID, name: str) -> Response:
    """Картинка/рисунок, извлечённый из оригинала (MinerU), — для вставки
    в MD-просмотр перевода. Ключ детерминированный: {doc_id}/img/{имя файла}."""
    await _get_or_404(request, doc_id)
    key = f"{doc_id}/img/{Path(name).name}"
    try:
        data = await request.app.state.storage.get_bytes(settings.bucket_artifacts, key)
    except Exception as exc:
        raise HTTPException(404, "картинка не найдена") from exc
    media = mimetypes.guess_type(name)[0] or "image/jpeg"
    return Response(content=data, media_type=media, headers={"Cache-Control": "public, max-age=86400"})


# --- интерактивный xlsx-просмотр (сетка листов, а не office-PDF «принт») ---
_SHEET_MAX_ROWS = 1000  # потолок строк на лист для грид-просмотра (латентность/payload)
_SHEET_MAX_COLS = 60


def _fmt_cell(v: object) -> str:
    """Значение ячейки → строка для грида: даты как ДД.ММ.ГГГГ, числа без
    «хвостов» float (1.8849999999999998 → 1.885), остальное — как есть."""
    import datetime as _dt

    if v is None:
        return ""
    if isinstance(v, _dt.datetime):
        return (
            v.strftime("%d.%m.%Y")
            if (v.hour, v.minute, v.second) == (0, 0, 0)
            else v.strftime("%d.%m.%Y %H:%M")
        )
    if isinstance(v, _dt.date):
        return v.strftime("%d.%m.%Y")
    if isinstance(v, bool):
        return "ИСТИНА" if v else "ЛОЖЬ"
    if isinstance(v, float):
        return f"{v:g}"
    return str(v)


def _ws_grids(ws: object, trans: dict[str, str]) -> tuple[list[list[str]], list[list[str]], int, int]:
    """Лист (data_only) → (сетка оригинала, сетка перевода, всего_строк, всего_столбцов).

    Перевод = та же сетка, где строковые ячейки с известным переводом заменены на
    перевод; числа/формулы (кэш-значения) остаются на месте."""
    tot_r = getattr(ws, "max_row", 0) or 0
    tot_c = getattr(ws, "max_column", 0) or 0
    n_r = min(tot_r, _SHEET_MAX_ROWS)
    n_c = min(tot_c, _SHEET_MAX_COLS)
    og: list[list[str]] = []
    rg: list[list[str]] = []
    if n_r and n_c:
        for row in ws.iter_rows(min_row=1, max_row=n_r, max_col=n_c, values_only=True):  # type: ignore[attr-defined]
            o_row: list[str] = []
            r_row: list[str] = []
            for v in row:
                cell = _fmt_cell(v)
                o_row.append(cell)
                t = trans.get(v) if isinstance(v, str) else None
                r_row.append(t if t else cell)
            og.append(o_row)
            rg.append(r_row)
    return og, rg, tot_r, tot_c


def _ws_chart_titles(ws: object) -> list[str]:
    """Заголовки встроенных диаграмм листа (грид показывает ячейки, не рисунки —
    о диаграмме сообщаем пометкой, иначе она «пропадает»)."""
    titles: list[str] = []
    for c in getattr(ws, "_charts", []) or []:
        t = getattr(c, "title", None)
        txt = None
        try:
            if t is not None and getattr(t, "tx", None) and getattr(t.tx, "rich", None):
                txt = t.tx.rich.p[0].r[0].t
        except Exception:
            txt = None
        titles.append(txt or "диаграмма")
    return titles


def _read_xlsx_grids(orig: bytes, trans: dict[str, str]) -> list[dict]:
    """Оригинал (data_only — с кэш-значениями формул) + словарь переводов
    {исходный_текст_ячейки: перевод} → листы с сетками orig/ru для грида.

    Перевод накладывается на оригинал из сегментов БД, а НЕ читается из
    переэкспортированного xlsx: openpyxl при сохранении теряет кэш формул, из-за
    чего числовые ячейки-формулы (=COUNTIF, =SUM, =IF) выглядели пустыми."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(orig), data_only=True)
    out: list[dict] = []
    for ws in wb.worksheets:
        og, rg, tot_r, tot_c = _ws_grids(ws, trans)
        out.append(
            {
                "name": ws.title,
                "name_ru": trans.get(ws.title) or ws.title,
                "orig": og,
                "ru": rg,
                "total_rows": tot_r,
                "total_cols": tot_c,
                "truncated": tot_r > _SHEET_MAX_ROWS or tot_c > _SHEET_MAX_COLS,
                "charts": _ws_chart_titles(ws),
            }
        )
    return out


@router.get("/{doc_id}/sheets")
async def document_sheets(request: Request, doc_id: uuid.UUID) -> dict:
    """Данные xlsx для ИНТЕРАКТИВНОГО просмотра: листы → сетка ячеек (оригинал +
    перевод). Перевод накладывается на оригинал по тексту ячейки (из сегментов),
    поэтому числа/формулы видны всегда. Это настоящая таблица, а не PDF-принт."""
    doc = await _get_or_404(request, doc_id)
    kind = doc.kind if isinstance(doc.kind, str) else doc.kind.value
    if kind != "xlsx":
        raise HTTPException(400, "интерактивная сетка только для xlsx")
    async with request.app.state.sessionmaker() as session:
        rows = (
            await session.execute(
                select(Segment.source_text, Segment.translated_text).where(
                    Segment.document_id == doc_id,
                    Segment.translated_text.is_not(None),
                )
            )
        ).all()
    trans = {src: tr for src, tr in rows if tr}
    storage = request.app.state.storage
    orig_bytes = await storage.get_bytes(settings.bucket_originals, doc.s3_key_original)
    sheets = await asyncio.to_thread(_read_xlsx_grids, orig_bytes, trans)
    return {"sheets": sheets, "translated_ready": bool(trans)}


# --- интерактивный pptx-просмотр (структура слайдов, а не office-PDF «принт») ---
def _slide_blocks(shapes: object, s_i: int, trans: dict[str, str]) -> list[dict]:
    """Фигуры слайда (с заходом в группы) → упорядоченные блоки для рендера:
    text (абзацы), table (ячейки), image (рисунок). Перевод накладывается по
    location-ключу из сегментов (как inject_pptx)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    blocks: list[dict] = []
    for shape in shapes:  # type: ignore[attr-defined]
        try:
            stype = shape.shape_type
        except Exception:
            stype = None
        if stype == MSO_SHAPE_TYPE.GROUP:
            blocks.extend(_slide_blocks(shape.shapes, s_i, trans))
            continue
        if getattr(shape, "has_table", False):
            rows: list[list[dict]] = []
            for r, row in enumerate(shape.table.rows):
                cells: list[dict] = []
                for c, cell in enumerate(row.cells):
                    o = (cell.text or "").strip()
                    key = ooxml.location_key({"slide": s_i, "shape": shape.shape_id, "row": r, "col": c})
                    cells.append({"orig": o, "ru": trans.get(key) or o})
                rows.append(cells)
            blocks.append({"type": "table", "rows": rows})
            continue
        if stype == MSO_SHAPE_TYPE.PICTURE:
            blocks.append({"type": "image", "shape": shape.shape_id})
            continue
        if getattr(shape, "has_text_frame", False):
            lines: list[dict] = []
            for p_i, p in enumerate(shape.text_frame.paragraphs):
                o = (p.text or "").strip()
                if not o:
                    continue
                key = ooxml.location_key({"slide": s_i, "shape": shape.shape_id, "para": p_i})
                lines.append({"orig": o, "ru": trans.get(key) or o, "level": getattr(p, "level", 0) or 0})
            if lines:
                blocks.append({"type": "text", "lines": lines})
    return blocks


def _build_slides(orig: bytes, trans: dict[str, str]) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(orig))
    out: list[dict] = []
    for s_i, slide in enumerate(prs.slides):
        title_o = ""
        title_r = ""
        try:
            ts = slide.shapes.title
        except Exception:
            ts = None
        if ts is not None and getattr(ts, "has_text_frame", False):
            title_o = (ts.text or "").strip()
            key = ooxml.location_key({"slide": s_i, "shape": ts.shape_id, "para": 0})
            title_r = trans.get(key) or title_o
        out.append(
            {
                "index": s_i,
                "title": title_o,
                "title_ru": title_r or title_o,
                "blocks": _slide_blocks(slide.shapes, s_i, trans),
            }
        )
    return out


@router.get("/{doc_id}/slides")
async def document_slides(request: Request, doc_id: uuid.UUID) -> dict:
    """Данные pptx для ИНТЕРАКТИВНОГО просмотра: слайды → блоки (текст/таблица/
    рисунок), оригинал + перевод. Перевод накладывается по location из сегментов.
    Настоящая презентация (листать слайды, выделять текст), а не PDF-принт."""
    doc = await _get_or_404(request, doc_id)
    kind = doc.kind if isinstance(doc.kind, str) else doc.kind.value
    if kind != "pptx":
        raise HTTPException(400, "структура слайдов только для pptx")
    async with request.app.state.sessionmaker() as session:
        segs = (
            await session.execute(
                select(Segment).where(
                    Segment.document_id == doc_id, Segment.translated_text.is_not(None)
                )
            )
        ).scalars().all()
    trans = {
        ooxml.location_key(s.meta["location"]): s.translated_text
        for s in segs
        if s.translated_text and s.meta.get("location")
    }
    storage = request.app.state.storage
    orig_bytes = await storage.get_bytes(settings.bucket_originals, doc.s3_key_original)
    slides = await asyncio.to_thread(_build_slides, orig_bytes, trans)
    return {"slides": slides, "translated_ready": bool(trans)}


def _extract_slide_image(orig: bytes, slide_idx: int, shape_id: int) -> tuple[bytes, str] | None:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(orig))
    slides = list(prs.slides)
    if slide_idx < 0 or slide_idx >= len(slides):
        return None

    def find(shapes: object):
        for shape in shapes:  # type: ignore[attr-defined]
            try:
                stype = shape.shape_type
            except Exception:
                stype = None
            if stype == MSO_SHAPE_TYPE.GROUP:
                hit = find(shape.shapes)
                if hit:
                    return hit
            elif shape.shape_id == shape_id and stype == MSO_SHAPE_TYPE.PICTURE:
                return shape
        return None

    shape = find(slides[slide_idx].shapes)
    if shape is None:
        return None
    img = shape.image
    return img.blob, (img.content_type or "image/png")


@router.get("/{doc_id}/slide-image/{slide_idx}/{shape_id}")
async def slide_image(request: Request, doc_id: uuid.UUID, slide_idx: int, shape_id: int) -> Response:
    """Рисунок со слайда pptx (по slide+shape_id) — для интерактивного просмотра."""
    doc = await _get_or_404(request, doc_id)
    orig_bytes = await request.app.state.storage.get_bytes(settings.bucket_originals, doc.s3_key_original)
    res = await asyncio.to_thread(_extract_slide_image, orig_bytes, slide_idx, shape_id)
    if res is None:
        raise HTTPException(404, "рисунок не найден")
    blob, media = res
    return Response(content=blob, media_type=media, headers={"Cache-Control": "public, max-age=86400"})


@router.delete("/{doc_id}", status_code=204)
async def delete_document(request: Request, doc_id: uuid.UUID) -> None:
    """Удаление документа из библиотеки: объекты в S3 (оригинал, артефакт парсинга,
    экспорты) + запись в БД (каскадом — сегменты, чанки, эмбеддинги страниц,
    чат-сессии этого документа) + чистка памяти, привязанной к документу."""
    doc = await _get_or_404(request, doc_id)
    storage = request.app.state.storage

    # 1) S3 — best-effort по всем известным ключам документа
    await storage.remove_object(settings.bucket_originals, doc.s3_key_original)
    if doc.s3_key_content_list:
        await storage.remove_object(settings.bucket_artifacts, doc.s3_key_content_list)
    for attr, _media in _EXPORT_KINDS.values():
        key = getattr(doc, attr)
        if key:
            await storage.remove_object(settings.bucket_exports, key)

    # 2) БД — удаление документа (FK ondelete=CASCADE уносит segments/chunks/
    #    page_embeddings/chat_sessions+messages этого документа)
    async with request.app.state.sessionmaker() as session:
        obj = await session.get(Document, doc_id)
        if obj is not None:
            await session.delete(obj)
            await session.commit()

    # 3) Память документа (слой памяти расцеплён с FK) — мягкое удаление; под
    #    RLS FORCE нужен GUC скоупа владельца. Никогда не блокирует удаление.
    user: User = request.state.user
    try:
        memory = request.app.state.memory
        async with request.app.state.sessionmaker() as session:
            await apply_scope_guc(session, memory.scope_for(user.sub, document_id=doc_id))
            await session.execute(
                sql(
                    "UPDATE memory_items SET status='deleted', deleted_at=now()"
                    " WHERE user_id=:u AND document_id=:d AND status<>'deleted'"
                ),
                {"u": user.sub, "d": str(doc_id)},
            )
            await session.commit()
    except Exception:  # noqa: BLE001 — чистка памяти не должна валить удаление
        pass

    await audit(request, "delete", "document", str(doc_id), {"filename": doc.filename})


async def _get_or_404(request: Request, doc_id: uuid.UUID) -> Document:
    async with request.app.state.sessionmaker() as session:
        doc = await session.get(Document, doc_id)
    if doc is None:
        raise HTTPException(404, "документ не найден")
    user: User = request.state.user
    if not user.is_admin and doc.owner_sub is not None and doc.owner_sub != user.sub:
        raise HTTPException(404, "документ не найден")  # не раскрываем существование
    return doc
