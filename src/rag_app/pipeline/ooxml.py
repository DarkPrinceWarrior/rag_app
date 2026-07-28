"""OOXML-транслятор (roadmap § 3.3.B): DOCX/XLSX/PPTX → тот же формат.

Вёрстку не трогаем — она в XML. Извлекаем текстовые узлы с адресом
(location в meta сегмента), переводим сегментно, записываем обратно
в копию оригинала.

Адресация:
- DOCX:  {"p": i} — абзац body; {"t": ti, "r": ri, "c": ci, "p": pi} — абзац ячейки
- XLSX:  {"sheet": name, "cell": "A1"}
- PPTX:  {"slide": si, "shape": id, "para": pi}; заметки — {"slide": si, "notes": true, "para": pi}
"""

from __future__ import annotations

import json
import logging
import re
from pathlib import Path
from typing import Any

from docx import Document as DocxDocument
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, Twips
from docx.table import Table
from docx.text.paragraph import Paragraph
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE

from rag_app.config import settings
from rag_app.db.models import SegmentKind
from rag_app.pipeline.segments import SegmentDraft

logger = logging.getLogger(__name__)


def location_key(location: dict[str, Any]) -> str:
    return json.dumps(location, sort_keys=True, ensure_ascii=False)


# ------------------------------------------------------------------ DOCX


def _docx_set_paragraph_text(paragraph: Any, text: str) -> None:
    """Перевод в первый run, включая очистку текста внутри гиперссылок."""
    if not paragraph.runs:
        # paragraph.runs не включает run'ы внутри w:hyperlink. Очищаем их явно,
        # иначе к переводу приписывается исходный текст ссылки.
        for text_node in paragraph._p.xpath(".//w:t"):
            text_node.text = ""
        if text:
            paragraph.add_run(text)
        return
    paragraph.runs[0].text = text
    first_text_nodes = set(paragraph.runs[0]._r.xpath(".//w:t"))
    for text_node in paragraph._p.xpath(".//w:t"):
        if text_node not in first_text_nodes:
            text_node.text = ""


def _docx_table_grid_widths(table: Any) -> list[int]:
    """Ширины gridCol в twip; пустой список для нестандартной/битой таблицы."""
    result: list[int] = []
    for column in table._tbl.xpath("./w:tblGrid/w:gridCol"):
        raw = column.get(qn("w:w"))
        try:
            result.append(int(raw))
        except (TypeError, ValueError):
            return []
    return result


def _adapt_sparse_translated_table(
    table: Any,
    original_rows: list[list[str]],
) -> bool:
    """Точечно адаптировать таблицу с широкой пустой колонкой под более объёмный RU.

    Реальный Annex выявил узкий, но воспроизводимый класс: ≥4 строк, одна широкая
    колонка занимает ≥45% таблицы и полностью пуста во всех строках данных, а
    перевод соседней содержательной колонки заметно длиннее оригинала. В
    исходной сетке пустая Conclusion
    занимала 68%, поэтому русский текст в Topics разрывал одну логическую строку
    между страницами.

    Для этого класса немного перераспределяем только существующую ширину
    (содержательной колонке до 32%, пустой остаётся не менее 45%), убираем
    исходные минимальные высоты строк, ставим 11 pt только содержательным
    абзацам данных и запрещаем межстраничный разрыв строки. Остальные таблицы
    остаются байт-в-байт по параметрам раскладки.
    """
    widths = _docx_table_grid_widths(table)
    if len(widths) < 3 or len(original_rows) < 4 or any(len(row) != len(widths) for row in original_rows):
        return False
    # Для merged cells одна и та же w:tc покрывает несколько gridCol. Простая
    # запись ширины отдельной колонки тогда перетирает суммарную ширину merge.
    if table._tbl.xpath(".//w:gridSpan") or table._tbl.xpath(".//w:vMerge"):
        return False
    # Строки фиксированной высоты характерны для бланков и мест ручного
    # заполнения. Их нельзя растягивать/сжимать эвристикой переноса текста.
    if any(
        height.get(qn("w:hRule")) == "exact"
        for row in table.rows
        for height in row._tr.xpath("./w:trPr/w:trHeight")
    ):
        return False
    total_width = sum(widths)
    if total_width <= 0:
        return False

    body_rows = original_rows[1:]
    donor_candidates = [
        ci
        for ci, width in enumerate(widths)
        if width / total_width >= 0.45 and not any(row[ci].strip() for row in body_rows)
    ]
    if not donor_candidates:
        return False
    donor = max(donor_candidates, key=widths.__getitem__)

    target_candidates: list[tuple[int, int]] = []
    for ci, width in enumerate(widths):
        if ci == donor or width / total_width >= 0.32:
            continue
        translated_chars = sum(
            len(row.cells[ci].text.strip()) for row in table.rows[1:] if ci < len(row.cells)
        )
        original_chars = sum(len(row[ci].strip()) for row in body_rows)
        changed = any(
            row.cells[ci].text.strip() != body_rows[ri][ci].strip() for ri, row in enumerate(table.rows[1:])
        )
        # Даже умеренное суммарное удлинение по нескольким строкам способно
        # вытолкнуть строку на следующую страницу (реальный Annex: +14 / 182).
        minimum_expansion = original_chars + max(8, original_chars // 20)
        if changed and original_chars >= 20 and translated_chars >= minimum_expansion:
            target_candidates.append((translated_chars, ci))
    if not target_candidates:
        return False
    target = max(target_candidates)[1]

    desired_target = round(total_width * 0.32)
    minimum_donor = round(total_width * 0.45)
    transfer = min(desired_target - widths[target], widths[donor] - minimum_donor)
    if transfer <= 0:
        return False
    widths[target] += transfer
    widths[donor] -= transfer

    grid_columns = table._tbl.xpath("./w:tblGrid/w:gridCol")
    for ci in (target, donor):
        grid_columns[ci].set(qn("w:w"), str(widths[ci]))
        seen_cells: set[int] = set()
        for row in table.rows:
            cell = row.cells[ci]
            identity = id(cell._tc)
            if identity in seen_cells:
                continue
            seen_cells.add(identity)
            cell.width = Twips(widths[ci])

    for row in table.rows:
        tr_properties = row._tr.get_or_add_trPr()
        for height in list(tr_properties.findall(qn("w:trHeight"))):
            if height.get(qn("w:hRule")) == "atLeast":
                tr_properties.remove(height)
        if tr_properties.find(qn("w:cantSplit")) is None:
            tr_properties.append(OxmlElement("w:cantSplit"))

    for row in table.rows[1:]:
        for paragraph in row.cells[target].paragraphs:
            if not paragraph.text.strip():
                continue
            for run in paragraph.runs:
                # Не увеличиваем шрифт документов, где он уже 11 pt или меньше;
                # у Annex 12 pt наследуются от paragraph style, а не заданы run.
                effective_size = run.font.size
                if effective_size is None and paragraph.style is not None:
                    effective_size = paragraph.style.font.size
                if run.text and effective_size is not None and effective_size > Pt(11):
                    run.font.size = Pt(11)
    return True


def _docx_paragraph_images(paragraph: Any, images_dir: Path, drafts: list[SegmentDraft]) -> None:
    """Встроенные картинки абзаца → файлы в images_dir + сегменты kind=image.

    Картинка лежит в part'ах документа, ссылка — `a:blip r:embed`. Кладём байты
    в images_dir (img_path в meta), парс-задача потом грузит их в MinIO для
    вставки в MD-просмотр. Сегмент-картинка идёт сразу за своим абзацем."""
    for blip in paragraph._p.iterfind(".//" + qn("a:blip")):
        rid = blip.get(qn("r:embed")) or blip.get(qn("r:link"))
        if not rid:
            continue
        try:
            part = paragraph.part.related_parts[rid]
        except KeyError:
            continue
        name = Path(str(part.partname)).name
        try:
            (images_dir / name).write_bytes(part.blob)
        except Exception:
            continue
        drafts.append(
            SegmentDraft(idx=len(drafts), kind=SegmentKind.image, source_text="", meta={"img_path": name})
        )


def extract_docx(path: Path, images_dir: Path | None = None) -> list[SegmentDraft]:
    doc = DocxDocument(str(path))
    drafts: list[SegmentDraft] = []

    def add(
        text: str,
        location: dict[str, Any],
        style_name: str,
        *,
        table_size: list[int] | None = None,
    ) -> None:
        text = text.strip()
        if not text:
            return
        kind = SegmentKind.paragraph
        level = None
        if style_name.startswith("Heading"):
            kind = SegmentKind.heading
            try:
                level = int(style_name.split()[-1])
            except ValueError:
                level = 1
        meta: dict[str, Any] = {"location": location}
        if table_size is not None:
            # Пустые ячейки не становятся переводимыми сегментами, поэтому без
            # размера таблицы viewer вычислял ширину только по непустым колонкам
            # и полностью терял, например, пустую колонку Conclusion.
            meta["table_size"] = table_size
        drafts.append(
            SegmentDraft(
                idx=len(drafts),
                kind=kind,
                source_text=text,
                heading_level=level,
                meta=meta,
            )
        )

    # идём по детям body В ПОРЯДКЕ ДОКУМЕНТА (абзацы, таблицы, картинки на своих
    # местах). Индексы p_idx/t_idx совпадают с doc.paragraphs[i]/doc.tables[ti],
    # поэтому location-ключи те же, что у inject_docx — экспорт не ломается.
    p_idx = 0
    t_idx = 0
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)
            add(para.text, {"p": p_idx}, para.style.name if para.style else "")
            if images_dir is not None:
                _docx_paragraph_images(para, images_dir, drafts)
            p_idx += 1
        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            table_size = [
                len(table.rows),
                max((len(row.cells) for row in table.rows), default=0),
            ]
            for ri, row in enumerate(table.rows):
                for ci, cell in enumerate(row.cells):
                    for pi, cp in enumerate(cell.paragraphs):
                        add(
                            cp.text,
                            {"t": t_idx, "r": ri, "c": ci, "p": pi},
                            "",
                            table_size=table_size,
                        )
            t_idx += 1
    return drafts


def inject_docx(src: Path, dst: Path, translations: dict[str, str]) -> int:
    doc = DocxDocument(str(src))
    original_table_rows = [[[cell.text for cell in row.cells] for row in table.rows] for table in doc.tables]
    applied = 0

    def apply(paragraph: Any, location: dict[str, Any]) -> None:
        nonlocal applied
        text = translations.get(location_key(location))
        if text is not None:
            _docx_set_paragraph_text(paragraph, text)
            applied += 1

    for i, p in enumerate(doc.paragraphs):
        apply(p, {"p": i})
    for ti, table in enumerate(doc.tables):
        for ri, row in enumerate(table.rows):
            for ci, cell in enumerate(row.cells):
                for pi, p in enumerate(cell.paragraphs):
                    apply(p, {"t": ti, "r": ri, "c": ci, "p": pi})
        if _adapt_sparse_translated_table(table, original_table_rows[ti]):
            logger.info(
                "DOCX table %d: адаптирована широкая пустая колонка для перевода",
                ti,
            )
    doc.save(str(dst))
    return applied


# ------------------------------------------------------------------ XLSX

# Слово = ≥2 подряд идущих буквы (латиница/кириллица). Чисто-числовой/кодовый
# дамп (0.43, 130/130/300, DMFA, pH, Eo) слов в этом смысле не даёт «прозы».
_WORD_RE = re.compile(r"[A-Za-zА-Яа-яЁё]{2,}")
_DEFAULT_SHEET_TITLE_RE = re.compile(r"Sheet\d*", re.IGNORECASE)


def _is_translatable_xlsx(v: str) -> bool:
    """Ячейка переводима, если в ней есть осмысленный текст (слово/фраза), а не
    голый код/число/идентификатор.

    Фразы и длинные подписи переводятся сразу. Для одиночного токена отличаем
    естественные слова от кодов: минимум четыре буквы; короткий ALL-CAPS — код,
    а обычный регистр (`Valve`, `Status`) — подпись. ALL-CAPS длиной от семи
    букв сохраняем как возможный заголовок (`PRESSURE`)."""
    s = v.strip()
    if not s or s.startswith("="):
        return False
    if not _WORD_RE.search(s):
        return False
    if " " in s or len(s) >= 12:
        return True
    token = s.strip(".,:;!?()[]{}")
    if len(token) < 4 or not token.isalpha():
        return False
    if token.isupper():
        return len(token) >= 7
    return token.islower() or token.istitle() or len(token) >= 6


def extract_xlsx(path: Path) -> list[SegmentDraft]:
    wb = load_workbook(str(path))  # data_only=False: формулы остаются формулами
    drafts: list[SegmentDraft] = []
    seen: set[str] = set()  # дедуп по исходному тексту: одинаковые строки = 1 сегмент
    capped = False
    skipped_dup = 0
    for s_i, ws in enumerate(wb.worksheets):
        # название листа — тоже переводим (вкладки листов показываем и на русском)
        title = (ws.title or "").strip()
        if (
            title
            and not _DEFAULT_SHEET_TITLE_RE.fullmatch(title)
            and _is_translatable_xlsx(title)
            and title not in seen
        ):
            seen.add(title)
            drafts.append(
                SegmentDraft(
                    idx=len(drafts),
                    kind=SegmentKind.paragraph,
                    source_text=title,
                    page_idx=s_i,
                    meta={"location": {"sheet": ws.title, "cell": "__sheet_title__"}, "sheet_title": True},
                )
            )
        for row in ws.iter_rows():
            for cell in row:
                v = cell.value
                # только строковые значения; формулы и числа не трогаем по построению
                if not isinstance(v, str) or not _is_translatable_xlsx(v):
                    continue
                if v in seen:
                    skipped_dup += 1
                    continue
                if len(drafts) >= settings.xlsx_max_segments:
                    capped = True
                    continue
                seen.add(v)
                drafts.append(
                    SegmentDraft(
                        idx=len(drafts),
                        kind=SegmentKind.paragraph,
                        source_text=v,
                        page_idx=s_i,
                        # location — первой встреченной ячейки с этим текстом; inject
                        # применяет перевод ПО ТЕКСТУ ко всем ячейкам-дубликатам.
                        meta={"location": {"sheet": ws.title, "cell": cell.coordinate}},
                    )
                )
    if capped:
        logger.warning(
            "extract_xlsx %s: достигнут потолок xlsx_max_segments=%d — часть прозовых "
            "ячеек отброшена (перевод неполный); скрытых дублей-ячеек: %d",
            path.name,
            settings.xlsx_max_segments,
            skipped_dup,
        )
    return drafts


def inject_xlsx(src: Path, dst: Path, translations: dict[str, str]) -> int:
    """Записать перевод обратно в .xlsx ПО ИСХОДНОМУ ТЕКСТУ ячейки.

    translations: {исходный_текст_ячейки: перевод}. Перевод применяется ко ВСЕМ
    ячейкам с тем же исходным текстом (дедуп на extract → один перевод
    раскладывается на все дубликаты)."""
    wb = load_workbook(str(src))
    applied = 0
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if not isinstance(cell.value, str):
                    continue
                text = translations.get(cell.value)
                if text is not None:
                    cell.value = text
                    applied += 1
    wb.save(str(dst))
    return applied


# ------------------------------------------------------------------ PPTX

_CITATION_RE = re.compile(r"^\s*\[\d+\]")  # элемент списка литературы: «[1] Stevenson…»


def is_pptx_citation(text: str) -> bool:
    """Запись библиографии вида «[1] …» — список литературы не переводим."""
    return bool(_CITATION_RE.match(text or ""))


def _para_text(p: Any) -> str:
    return "".join(run.text for run in p.runs)


def _set_para(p: Any, text: str) -> None:
    """Записать перевод в абзац, сохранив форматирование первого run'а."""
    if p.runs:
        p.runs[0].text = text
        for run in p.runs[1:]:
            run.text = ""
    elif text:
        p.text = text


def _iter_shape_units(shapes: Any, s_i: int):
    """Рекурсивный обход фигур слайда (С ЗАХОДОМ В ГРУППЫ) → текстовые единицы.

    Yields (location, get_text, set_text) для:
    - абзацев текстовых фреймов  → {"slide", "shape", "para"}
    - ячеек таблиц               → {"slide", "shape", "row", "col"}
    Группы (MSO GROUP) рекурсивно разворачиваются — иначе текст в сгруппированных
    фигурах теряется (слайды-«только заголовок»). Таблицы (GraphicFrame.has_table)
    раньше вообще не извлекались."""
    for shape in shapes:
        try:
            is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            yield from _iter_shape_units(shape.shapes, s_i)
            continue
        if getattr(shape, "has_table", False):
            tbl = shape.table
            for r, row in enumerate(tbl.rows):
                for c, cell in enumerate(row.cells):
                    loc = {"slide": s_i, "shape": shape.shape_id, "row": r, "col": c}

                    def _get_cell(cell=cell) -> str:
                        return cell.text

                    def _set_cell(text: str, cell=cell) -> None:
                        cell.text = text

                    yield loc, _get_cell, _set_cell
            continue
        if getattr(shape, "has_text_frame", False):
            for p_i, p in enumerate(shape.text_frame.paragraphs):
                loc = {"slide": s_i, "shape": shape.shape_id, "para": p_i}

                def _get_paragraph(p=p) -> str:
                    return _para_text(p)

                def _set_paragraph(text: str, p=p) -> None:
                    _set_para(p, text)

                yield loc, _get_paragraph, _set_paragraph


def _pptx_units(prs: Any):
    """Все переводимые единицы презентации (фигуры+группы+таблицы и заметки)."""
    for s_i, slide in enumerate(prs.slides):
        yield from _iter_shape_units(slide.shapes, s_i)
        if slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            for p_i, p in enumerate(tf.paragraphs):
                loc = {"slide": s_i, "notes": True, "para": p_i}

                def _get(p=p) -> str:
                    return _para_text(p)

                def _set(text: str, p=p) -> None:
                    _set_para(p, text)

                yield loc, _get, _set


def extract_pptx(path: Path) -> list[SegmentDraft]:
    prs = Presentation(str(path))
    drafts: list[SegmentDraft] = []
    for location, get_text, _ in _pptx_units(prs):
        text = (get_text() or "").strip()
        if not text or is_pptx_citation(text):  # список литературы не переводим
            continue
        drafts.append(
            SegmentDraft(
                idx=len(drafts),
                kind=SegmentKind.paragraph,
                source_text=text,
                page_idx=location["slide"],
                meta={"location": location},
            )
        )
    return drafts


def inject_pptx(src: Path, dst: Path, translations: dict[str, str]) -> int:
    prs = Presentation(str(src))
    applied = 0
    for location, _, set_text in _pptx_units(prs):
        text = translations.get(location_key(location))
        if text is None:
            continue
        set_text(text)
        applied += 1
    prs.save(str(dst))
    return applied


_CM = 360000  # EMU в сантиметре


def _iter_geom_shapes(shapes: Any):
    """Плоский обход фигур (с заходом в группы) — для геометрии/автоподгонки."""
    for shape in shapes:
        try:
            is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:
            is_group = False
        if is_group:
            yield from _iter_geom_shapes(shape.shapes)
        else:
            yield shape


def _shrink_table(shape: Any) -> None:
    """Ужать таблицу: меньше кегль + минимальные поля ячеек, чтобы переведённый
    (более длинный) текст переносился на меньшее число строк и таблица не уезжала
    за нижний край слайда."""
    from pptx.util import Pt

    for row in shape.table.rows:
        for cell in row.cells:
            tf = cell.text_frame
            try:
                tf.word_wrap = True
                cell.margin_top = Pt(1)
                cell.margin_bottom = Pt(1)
                cell.margin_left = Pt(2)
                cell.margin_right = Pt(2)
            except Exception:
                pass
            for p in tf.paragraphs:
                for run in p.runs:
                    try:
                        sz = run.font.size
                        run.font.size = max(Pt(7), int(sz * 0.62)) if sz is not None else Pt(8)
                    except Exception:
                        pass


def _autofit_slide(slide: Any) -> None:
    # препятствия (картинки/таблицы) — их верхняя кромка; текстовый блок, который
    # их перекрывает, обрежем по высоте до них, чтобы текст не наезжал на картинку.
    obstacles: list[tuple[int, int]] = []  # (top, shape_id)
    for sh in _iter_geom_shapes(slide.shapes):
        try:
            is_pic = sh.shape_type == MSO_SHAPE_TYPE.PICTURE
        except Exception:
            is_pic = False
        if (getattr(sh, "has_table", False) or is_pic) and sh.top is not None and sh.height is not None:
            obstacles.append((int(sh.top), sh.shape_id))

    for sh in _iter_geom_shapes(slide.shapes):
        if getattr(sh, "has_table", False):
            _shrink_table(sh)
            continue
        try:
            is_pic = sh.shape_type == MSO_SHAPE_TYPE.PICTURE
        except Exception:
            is_pic = False
        if is_pic or not getattr(sh, "has_text_frame", False):
            continue
        tf = sh.text_frame
        # обрезать высоту блока до ближайшего препятствия ниже его верхней кромки
        if sh.top is not None and sh.height is not None:
            top = int(sh.top)
            below = [
                ot
                for (ot, oid) in obstacles
                if oid != sh.shape_id and top + _CM // 3 < ot < top + int(sh.height)
            ]
            if below:
                new_h = min(below) - top - _CM // 7  # небольшой зазор
                if new_h > _CM // 2:
                    try:
                        sh.height = int(new_h)
                    except Exception:
                        pass
        try:
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass


def pptx_autofit(src: Path, dst: Path) -> None:
    """Копия pptx для РЕНДЕРА-ПРОСМОТРА (office-PDF), подогнанная под фиксированный
    размер слайда: перенос слов + «ужать текст до фигуры», обрезка текстовых
    блоков, перекрывающих картинки, и уменьшение кегля таблиц. Нужна потому, что
    переведённый (русский) текст длиннее английского и в исходной раскладке
    наезжает на картинки / выходит за нижний край. Оригинальный .pptx-экспорт
    не трогаем."""
    prs = Presentation(str(src))
    for slide in prs.slides:
        _autofit_slide(slide)
    prs.save(str(dst))


# ------------------------------------------------------------------ единый вход

EXTRACTORS = {"docx": extract_docx, "xlsx": extract_xlsx, "pptx": extract_pptx}
INJECTORS = {"docx": inject_docx, "xlsx": inject_xlsx, "pptx": inject_pptx}


def extract(kind: str, path: Path, images_dir: Path | None = None) -> list[SegmentDraft]:
    if kind == "docx":
        return extract_docx(path, images_dir)
    return EXTRACTORS[kind](path)


def inject(kind: str, src: Path, dst: Path, translations: dict[str, str]) -> int:
    return INJECTORS[kind](src, dst, translations)
