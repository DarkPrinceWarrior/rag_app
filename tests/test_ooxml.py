from __future__ import annotations

from pathlib import Path

from docx import Document as DocxDocument
from docx.enum.table import WD_ROW_HEIGHT_RULE
from docx.opc.constants import RELATIONSHIP_TYPE as RT
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Pt, Twips
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches

from rag_app.db.models import SegmentKind
from rag_app.pipeline import ooxml


def _add_hyperlink(paragraph, text: str, url: str) -> None:
    relationship_id = paragraph.part.relate_to(url, RT.HYPERLINK, is_external=True)
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), relationship_id)
    run = OxmlElement("w:r")
    text_node = OxmlElement("w:t")
    text_node.text = text
    run.append(text_node)
    hyperlink.append(run)
    paragraph._p.append(hyperlink)


def test_docx_roundtrip(tmp_path: Path) -> None:
    src = tmp_path / "src.docx"
    doc = DocxDocument()
    doc.add_heading("Scope", level=1)
    p = doc.add_paragraph()
    p.add_run("Design pressure is ")
    bold = p.add_run("16.5 MPa")
    bold.bold = True
    table = doc.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Item"
    table.cell(0, 1).text = "Value"
    doc.save(str(src))

    drafts = ooxml.extract_docx(src)
    texts = [d.source_text for d in drafts]
    assert "Scope" in texts
    assert "Design pressure is 16.5 MPa" in texts
    assert "Item" in texts
    heading = next(d for d in drafts if d.source_text == "Scope")
    assert heading.kind == SegmentKind.heading and heading.heading_level == 1
    table_drafts = [d for d in drafts if d.meta["location"].get("t") == 0]
    assert table_drafts
    assert all(d.meta["table_size"] == [1, 2] for d in table_drafts)

    translations = {ooxml.location_key(d.meta["location"]): f"RU:{d.source_text}" for d in drafts}
    dst = tmp_path / "dst.docx"
    applied = ooxml.inject_docx(src, dst, translations)
    assert applied == len(drafts)

    out = DocxDocument(str(dst))
    out_texts = [p.text for p in out.paragraphs if p.text.strip()]
    assert "RU:Scope" in out_texts
    assert "RU:Design pressure is 16.5 MPa" in out_texts
    assert out.tables[0].cell(0, 0).text == "RU:Item"


def test_docx_extraction_keeps_empty_table_columns_in_metadata(tmp_path: Path) -> None:
    src = tmp_path / "empty-column.docx"
    doc = DocxDocument()
    table = doc.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "No."
    table.cell(0, 1).text = "Topic"
    table.cell(0, 2).text = "Conclusion"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "Design review"
    # Последняя ячейка намеренно пустая, как в реальном Annex.
    doc.save(src)

    drafts = ooxml.extract_docx(src)
    table_drafts = [draft for draft in drafts if "t" in draft.meta["location"]]

    assert len(table_drafts) == 5
    assert all(draft.meta["table_size"] == [2, 3] for draft in table_drafts)


def test_docx_injection_adapts_sparse_wide_table_without_splitting_rows(
    tmp_path: Path,
) -> None:
    src = tmp_path / "annex-like.docx"
    doc = DocxDocument()
    doc.styles["Normal"].font.size = Pt(12)
    table = doc.add_table(rows=5, cols=3)
    table.autofit = False
    widths = [588, 2034, 5615]
    for ci, width in enumerate(widths):
        table.columns[ci].width = Twips(width)
        for row in table.rows:
            row.cells[ci].width = Twips(width)
    headers = ["Sn.", "Meeting topics", "Conclusion"]
    for ci, text in enumerate(headers):
        table.cell(0, ci).text = text
    for ri in range(1, 5):
        table.cell(ri, 0).text = str(ri)
        table.cell(ri, 1).text = f"Engineering topic {ri} and review plan"
        table.cell(ri, 2).text = ""
        table.rows[ri].height = Twips(1400 + ri * 100)
        table.rows[ri].height_rule = WD_ROW_HEIGHT_RULE.AT_LEAST
    doc.save(src)

    drafts = ooxml.extract_docx(src)
    translations = {
        ooxml.location_key(draft.meta["location"]): (
            f"Инженерная тема {draft.meta['location']['r']} и подробный план проверки"
            if draft.meta["location"].get("c") == 1 and draft.meta["location"].get("r", 0) > 0
            else draft.source_text
        )
        for draft in drafts
    }
    dst = tmp_path / "annex-like.ru.docx"
    ooxml.inject_docx(src, dst, translations)

    original = DocxDocument(src)
    translated = DocxDocument(dst)
    original_widths = ooxml._docx_table_grid_widths(original.tables[0])
    translated_widths = ooxml._docx_table_grid_widths(translated.tables[0])

    assert sum(translated_widths) == sum(original_widths)
    minimum_donor = round(sum(original_widths) * 0.45)
    assert translated_widths[2] == minimum_donor
    assert translated_widths[1] == original_widths[1] + original_widths[2] - minimum_donor
    for row in translated.tables[0].rows:
        properties = row._tr.get_or_add_trPr()
        assert properties.find(qn("w:trHeight")) is None
        assert properties.find(qn("w:cantSplit")) is not None
        for ci in (1, 2):
            tc_width = row.cells[ci]._tc.get_or_add_tcPr().get_or_add_tcW().get(qn("w:w"))
            assert int(tc_width) == translated_widths[ci]
    first_topic_run = translated.tables[0].cell(1, 1).paragraphs[0].runs[0]
    assert first_topic_run.font.size == Pt(11)


def test_docx_sparse_table_adaptation_skips_merged_cells(tmp_path: Path) -> None:
    src = tmp_path / "merged-header.docx"
    doc = DocxDocument()
    table = doc.add_table(rows=5, cols=3)
    table.autofit = False
    widths = [588, 2034, 5615]
    for ci, width in enumerate(widths):
        table.columns[ci].width = Twips(width)
        for row in table.rows:
            row.cells[ci].width = Twips(width)
    table.cell(0, 0).text = "Sn."
    merged = table.cell(0, 1).merge(table.cell(0, 2))
    merged.text = "Meeting topics and conclusion"
    for ri in range(1, 5):
        table.cell(ri, 0).text = str(ri)
        table.cell(ri, 1).text = f"Engineering topic {ri} and review plan"
        table.cell(ri, 2).text = ""
    doc.save(src)

    original = DocxDocument(src)
    original_widths = ooxml._docx_table_grid_widths(original.tables[0])
    original_merged_width = original.tables[0].cell(0, 1)._tc.xpath("./w:tcPr/w:tcW")[0].get(qn("w:w"))

    drafts = ooxml.extract_docx(src)
    translations = {
        ooxml.location_key(draft.meta["location"]): (
            "Подробная инженерная тема и расширенный план совместной проверки"
            if draft.meta["location"].get("r", 0) > 0 and draft.meta["location"].get("c") == 1
            else draft.source_text
        )
        for draft in drafts
    }
    dst = tmp_path / "merged-header.ru.docx"
    ooxml.inject_docx(src, dst, translations)

    translated = DocxDocument(dst)
    assert ooxml._docx_table_grid_widths(translated.tables[0]) == original_widths
    translated_merged_width = translated.tables[0].cell(0, 1)._tc.xpath("./w:tcPr/w:tcW")[0].get(qn("w:w"))
    assert translated_merged_width == original_merged_width


def test_docx_sparse_table_adaptation_does_not_expand_small_font_or_short_translation(
    tmp_path: Path,
) -> None:
    src = tmp_path / "short-translation.docx"
    doc = DocxDocument()
    table = doc.add_table(rows=5, cols=3)
    table.autofit = False
    widths = [588, 2034, 5615]
    for ci, width in enumerate(widths):
        table.columns[ci].width = Twips(width)
        for row in table.rows:
            row.cells[ci].width = Twips(width)
    for ci, text in enumerate(["Sn.", "Meeting topics", "Conclusion"]):
        table.cell(0, ci).text = text
    for ri in range(1, 5):
        table.cell(ri, 0).text = str(ri)
        table.cell(ri, 1).text = f"Very long engineering coordination topic number {ri}"
        table.cell(ri, 1).paragraphs[0].runs[0].font.size = Pt(8)
        table.rows[ri].height = Twips(1200)
        table.rows[ri].height_rule = WD_ROW_HEIGHT_RULE.AT_LEAST
    doc.save(src)

    drafts = ooxml.extract_docx(src)
    translations = {
        ooxml.location_key(draft.meta["location"]): (
            f"Тема {draft.meta['location']['r']}"
            if draft.meta["location"].get("r", 0) > 0 and draft.meta["location"].get("c") == 1
            else draft.source_text
        )
        for draft in drafts
    }
    dst = tmp_path / "short-translation.ru.docx"
    ooxml.inject_docx(src, dst, translations)

    translated = DocxDocument(dst)
    assert ooxml._docx_table_grid_widths(translated.tables[0]) == widths
    first_body_row = translated.tables[0].rows[1]
    assert first_body_row.height == Twips(1200)
    assert first_body_row.cells[1].paragraphs[0].runs[0].font.size == Pt(8)


def test_docx_sparse_table_adaptation_preserves_exact_row_heights(
    tmp_path: Path,
) -> None:
    src = tmp_path / "fixed-height-form.docx"
    doc = DocxDocument()
    doc.styles["Normal"].font.size = Pt(12)
    table = doc.add_table(rows=5, cols=3)
    table.autofit = False
    widths = [588, 2034, 5615]
    for ci, width in enumerate(widths):
        table.columns[ci].width = Twips(width)
        for row in table.rows:
            row.cells[ci].width = Twips(width)
    for ci, text in enumerate(["No.", "Requirement", "Handwritten notes"]):
        table.cell(0, ci).text = text
    for ri in range(1, 5):
        table.cell(ri, 0).text = str(ri)
        table.cell(ri, 1).text = f"Engineering inspection requirement {ri}"
        table.rows[ri].height = Twips(1800)
        table.rows[ri].height_rule = WD_ROW_HEIGHT_RULE.EXACTLY
    doc.save(src)

    drafts = ooxml.extract_docx(src)
    translations = {
        ooxml.location_key(draft.meta["location"]): (
            "Подробное инженерное требование к совместной технической проверке"
            if draft.meta["location"].get("r", 0) > 0 and draft.meta["location"].get("c") == 1
            else draft.source_text
        )
        for draft in drafts
    }
    dst = tmp_path / "fixed-height-form.ru.docx"
    ooxml.inject_docx(src, dst, translations)

    translated = DocxDocument(dst)
    assert ooxml._docx_table_grid_widths(translated.tables[0]) == widths
    for row in translated.tables[0].rows[1:]:
        assert row.height == Twips(1800)
        assert row.height_rule == WD_ROW_HEIGHT_RULE.EXACTLY


def test_docx_injection_does_not_duplicate_hyperlink_text(tmp_path: Path) -> None:
    src = tmp_path / "hyperlink.docx"
    doc = DocxDocument()
    paragraph = doc.add_paragraph("See spec at ")
    _add_hyperlink(paragraph, "example.com/spec", "https://example.com/spec")
    only_link = doc.add_paragraph()
    _add_hyperlink(only_link, "original link", "https://example.com")
    doc.save(src)

    drafts = ooxml.extract_docx(src)
    assert [draft.source_text for draft in drafts] == [
        "See spec at example.com/spec",
        "original link",
    ]
    translations = {
        ooxml.location_key(drafts[0].meta["location"]): "См. спецификацию",
        ooxml.location_key(drafts[1].meta["location"]): "Переведённая ссылка",
    }
    dst = tmp_path / "translated.docx"
    assert ooxml.inject_docx(src, dst, translations) == 2

    out = DocxDocument(dst)
    assert [paragraph.text for paragraph in out.paragraphs] == [
        "См. спецификацию",
        "Переведённая ссылка",
    ]


def test_xlsx_roundtrip(tmp_path: Path) -> None:
    src = tmp_path / "src.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Design pressure"
    ws["B1"] = 16.5  # число — не трогаем
    ws["C1"] = "=B1*2"  # формула — не трогаем
    ws["A2"] = "Test pressure"
    ws["A3"] = "Design pressure"  # дубль текста A1 → один сегмент, оба перевода
    wb.save(str(src))

    drafts = ooxml.extract_xlsx(src)
    # дедуп по тексту: "Design pressure" один раз
    assert sorted(d.source_text for d in drafts) == ["Design pressure", "Test pressure"]

    # inject теперь по ИСХОДНОМУ ТЕКСТУ ячейки (а не по location)
    translations = {d.source_text: f"RU:{d.source_text}" for d in drafts}
    dst = tmp_path / "dst.xlsx"
    # 3 ячейки переведены (A1, A2, A3 — дубль тоже), несмотря на 2 сегмента
    assert ooxml.inject_xlsx(src, dst, translations) == 3

    out = load_workbook(str(dst))
    ws2 = out.active
    assert ws2["A1"].value == "RU:Design pressure"
    assert ws2["A3"].value == "RU:Design pressure"  # дубль получил тот же перевод
    assert ws2["B1"].value == 16.5
    assert ws2["C1"].value == "=B1*2"


def test_xlsx_skips_data_dump(tmp_path: Path) -> None:
    """Числовой/кодовый дамп НЕ плодит сегменты; проза — переводится."""
    src = tmp_path / "dump.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Compound name"  # проза (есть пробел) → переводим
    ws["B1"] = "Melting point, C"  # проза → переводим
    ws["A2"] = "DMFA"  # короткий код-токен → НЕ переводим
    ws["B2"] = "0.43"  # число-строка → НЕ переводим
    ws["A3"] = "130/130/300"  # код → НЕ переводим
    ws["B3"] = "pH"  # короткий токен → НЕ переводим
    ws["A4"] = "Eo"  # короткий токен → НЕ переводим
    ws["B4"] = "BTC"  # короткий код-токен → НЕ переводим
    ws["A5"] = "Hydrostatic test"  # проза → переводим
    wb.save(str(src))

    texts = sorted(d.source_text for d in ooxml.extract_xlsx(src))
    assert texts == ["Compound name", "Hydrostatic test", "Melting point, C"]


def test_xlsx_keeps_meaningful_single_word_labels_and_sheet_title(tmp_path: Path) -> None:
    src = tmp_path / "labels.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Process Data"
    ws.append(["Valve", "Status", "Pressure", "PRESSURE"])
    wb.save(str(src))

    drafts = ooxml.extract_xlsx(src)
    texts = {draft.source_text for draft in drafts}

    assert texts == {"Process Data", "Valve", "Status", "Pressure", "PRESSURE"}
    title = next(draft for draft in drafts if draft.source_text == "Process Data")
    assert title.meta["sheet_title"] is True


def test_pptx_roundtrip(tmp_path: Path) -> None:
    src = tmp_path / "src.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hydrostatic testing"
    prs.save(str(src))

    drafts = ooxml.extract_pptx(src)
    texts = [d.source_text for d in drafts]
    assert "Hydrostatic testing" in texts

    translations = {ooxml.location_key(d.meta["location"]): f"RU:{d.source_text}" for d in drafts}
    dst = tmp_path / "dst.pptx"
    applied = ooxml.inject_pptx(src, dst, translations)
    assert applied == len(drafts)

    out = Presentation(str(dst))
    out_texts = [
        sh.text_frame.text for s in out.slides for sh in s.shapes if getattr(sh, "has_text_frame", False)
    ]
    assert any("RU:Hydrostatic testing" in t for t in out_texts)


def test_pick_glossary_terms() -> None:
    from rag_app.llm.client import pick_glossary_terms

    terms = [
        ("maximum allowable working pressure", "максимально допустимое рабочее давление"),
        ("pressure vessel", "сосуд под давлением"),
        ("weld", "сварной шов"),
    ]
    found = pick_glossary_terms("The Pressure Vessel shall be welded.", terms)
    assert ("pressure vessel", "сосуд под давлением") in found
    # "weld" не должен находиться внутри "welded" (границы слов)
    assert ("weld", "сварной шов") not in found
